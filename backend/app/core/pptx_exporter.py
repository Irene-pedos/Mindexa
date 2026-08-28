from __future__ import annotations

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.agents.slide_deck_agent import SlideDeckOutput


def build_slide_deck_pptx(deck: SlideDeckOutput) -> bytes:
    """
    Render a validated SlideDeckOutput schema into a clean, professional .pptx presentation file.
    """
    prs = Presentation()
    
    # 16:9 widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide (Layout 0)
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    if title_slide.shapes.title:
        title_slide.shapes.title.text = deck.title
        title_para = title_slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(15, 23, 42) # Slate 900
    
    placeholders = title_slide.placeholders
    if len(placeholders) > 1:
        subtitle_shape = placeholders[1]
        subtitle_shape.text = (
            f"Audience: {deck.target_audience} • Duration: ~{deck.estimated_minutes} min\n"
            "Mindexa Academic Operating System"
        )
        sub_para = subtitle_shape.text_frame.paragraphs[0]
        sub_para.font.size = Pt(18)
        sub_para.font.color.rgb = RGBColor(100, 116, 139) # Slate 500

    # 2. Content Slides (Layout 1)
    content_layout = prs.slide_layouts[1]
    for slide_data in deck.slides:
        slide = prs.slides.add_slide(content_layout)
        
        # Slide Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
            title_para = slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(15, 23, 42)

        # Slide Body / Bullets
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for idx, bullet in enumerate(slide_data.bullet_points):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(30, 41, 59)
                p.space_after = Pt(14)

            # Optional visual cue note
            if slide_data.visual_idea:
                p_vis = tf.add_paragraph()
                p_vis.text = f"[Visual Cue: {slide_data.visual_idea}]"
                p_vis.font.size = Pt(14)
                p_vis.font.italic = True
                p_vis.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
                p_vis.space_before = Pt(16)

        # Speaker Notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = slide_data.speaker_notes

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
