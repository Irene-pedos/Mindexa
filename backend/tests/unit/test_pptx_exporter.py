import io
from pptx import Presentation

from app.agents.slide_deck_agent import SlideDeckOutput, SlideItem
from app.core.pptx_exporter import build_slide_deck_pptx


def test_build_slide_deck_pptx_generates_valid_presentation():
    deck = SlideDeckOutput(
        title="Introduction to Database Systems",
        target_audience="First Year CS",
        estimated_minutes=50,
        slides=[
            SlideItem(
                title="What is a DBMS?",
                bullet_points=["Software managing structured data", "ACID compliance guarantees", "Multi-user concurrency"],
                visual_idea="DBMS architecture box diagram",
                speaker_notes="Begin by asking students how they currently store state in applications."
            ),
            SlideItem(
                title="Relational Model Overview",
                bullet_points=["Tables, rows, and columns", "Primary keys and Foreign keys", "Declarative SQL interface"],
                visual_idea="ER diagram sample",
                speaker_notes="Emphasize mathematical foundations by Codd."
            ),
            SlideItem(
                title="ACID Properties",
                bullet_points=["Atomicity (all or nothing)", "Consistency (valid states)", "Isolation (concurrency)", "Durability (persistence)"],
                visual_idea="Transaction commit timeline",
                speaker_notes="Give a classic bank transfer example."
            ),
            SlideItem(
                title="Summary & Questions",
                bullet_points=["Key takeaways review", "Lab assignment 1 release", "Office hours schedule"],
                speaker_notes="Open the floor for questions."
            ),
        ]
    )

    pptx_bytes = build_slide_deck_pptx(deck)

    assert isinstance(pptx_bytes, bytes)
    assert len(pptx_bytes) > 0

    # Verify that python-pptx can parse the generated file
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    # 1 Title Slide + 4 Content Slides = 5 Slides
    assert len(prs.slides) == 5

    # Check title slide
    title_slide = prs.slides[0]
    assert title_slide.shapes.title.text == "Introduction to Database Systems"

    # Check content slide 1
    content_slide_1 = prs.slides[1]
    assert content_slide_1.shapes.title.text == "What is a DBMS?"

    # Check speaker notes on slide 1
    assert content_slide_1.notes_slide.notes_text_frame.text == "Begin by asking students how they currently store state in applications."
